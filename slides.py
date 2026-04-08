# ─────────────────────────────────────────────────────────────────────────────
# slides.py — PPTX deck generator from model data (no Streamlit dependency)
# ─────────────────────────────────────────────────────────────────────────────
import io
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE

# ─── Color palette ───────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x8A)
BLUE_MED   = RGBColor(0x1A, 0x56, 0xDB)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY       = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
RED        = RGBColor(0xDC, 0x26, 0x26)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_metric_card(slide, left, top, width, height, label, value,
                     positive=None):
    if positive is True:
        bg_color = RGBColor(0xD4, 0xED, 0xDA)
        val_color = GREEN
    elif positive is False:
        bg_color = RGBColor(0xF8, 0xD7, 0xDA)
        val_color = RED
    else:
        bg_color = BLUE_LIGHT
        val_color = BLUE_MED

    shape = _add_shape(slide, left, top, width, height, bg_color)
    shape.shadow.inherit = False

    # Label
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                  width - Inches(0.3), Inches(0.3),
                  label, font_size=9, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    # Value
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.35),
                  width - Inches(0.3), Inches(0.4),
                  value, font_size=18, color=val_color, bold=True, align=PP_ALIGN.CENTER)


def _add_table(slide, left, top, width, rows, cols, data, col_widths=None):
    """Add a styled table.
    data: list of lists (rows x cols), first row = headers.
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width,
                                          Inches(0.35 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx in range(rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(data[r_idx][c_idx])
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.name = "Calibri"

            if r_idx == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE_MED
                p.font.color.rgb = WHITE
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
            else:
                if c_idx == 0:
                    p.font.bold = True
                    p.alignment = PP_ALIGN.LEFT
                else:
                    p.alignment = PP_ALIGN.RIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else GRAY_LIGHT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def _add_chart_image(slide, fig, left, top, width, height):
    """Render a Plotly figure to PNG and embed in slide."""
    img_bytes = fig.to_image(format="png", width=int(width / Emu(1) * 96 / 914400),
                              height=int(height / Emu(1) * 96 / 914400), scale=2)
    stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(stream, left, top, width, height)


def _slide_header(slide, title, subtitle=None):
    """Standard slide header bar."""
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), BLUE_DARK)
    _add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.5),
                  title, font_size=24, bold=True, color=WHITE)
    if subtitle:
        _add_text_box(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.35),
                      subtitle, font_size=11, color=RGBColor(0xBF, 0xDB, 0xFE))


# ─── Slide builders ──────────────────────────────────────────────────────────

def _slide_cover(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, BLUE_DARK)

    # Accent bar
    _add_shape(slide, Inches(0.8), Inches(2.5), Inches(0.08), Inches(2.0), BLUE_MED)

    # Title
    _add_text_box(slide, Inches(1.2), Inches(2.5), Inches(9), Inches(0.8),
                  ctx["nome_proj"] or "Business Case", font_size=36, bold=True, color=WHITE)
    # Subtitle
    _add_text_box(slide, Inches(1.2), Inches(3.3), Inches(9), Inches(0.5),
                  ctx["L"]["slide_cover_subtitle"], font_size=16, color=RGBColor(0x93, 0xC5, 0xFD))

    # Metadata
    meta = []
    if ctx["setor"]:
        meta.append(f"{ctx['L']['slide_sector']}: {ctx['setor']}")
    if ctx["responsavel"]:
        meta.append(f"{ctx['L']['slide_author']}: {ctx['responsavel']}")
    if ctx["data_inicio"]:
        meta.append(f"{ctx['L']['slide_date']}: {ctx['data_inicio']}")
    _add_text_box(slide, Inches(1.2), Inches(4.0), Inches(9), Inches(0.6),
                  "  |  ".join(meta), font_size=11, color=RGBColor(0x93, 0xC5, 0xFD))

    # Footer
    _add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
                  "Business Case Analyzer v3.0", font_size=9, color=RGBColor(0x60, 0xA5, 0xFA))


def _slide_executive_summary(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_exec_title"], ctx["L"]["slide_exec_sub"])

    L = ctx["L"]
    # Verdict box
    npv = ctx["npv"]
    payback = ctx["payback"]
    horizonte = ctx["horizonte"]
    lim = horizonte * (2 / 3)

    if payback and payback <= lim and npv > 0:
        v_color = GREEN
        v_text = L["slide_viable"]
    elif (payback and payback <= horizonte) or npv > 0:
        v_color = ORANGE
        v_text = L["slide_caution"]
    else:
        v_color = RED
        v_text = L["slide_not_viable"]

    vbox = _add_shape(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.6), v_color)
    _add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.5),
                  v_text, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Metric cards — Row 1 (project)
    y = Inches(2.2)
    card_w = Inches(2.8)
    card_h = Inches(0.9)
    gap = Inches(0.2)
    x0 = Inches(0.6)

    fv = ctx["fv"]
    unit = ctx["unit"]

    cards_r1 = [
        (L["slide_roi"], f"{ctx['roi']:.0f}%", ctx['roi'] > 0),
        (L["slide_npv"], fv(npv, unit), npv > 0),
        (L["slide_payback"], f"{payback} {L['slide_months']}" if payback else "N/A",
         payback is not None),
        (L["slide_margin"], f"{ctx['mg_med']:.0f}%", None),
    ]
    for i, (lbl, val, pos) in enumerate(cards_r1):
        _add_metric_card(slide, x0 + i * (card_w + gap), y, card_w, card_h, lbl, val, pos)

    # Row 2 — IRR, MIRR, WACC, PI
    y2 = Inches(3.3)
    irr = ctx.get("irr_projeto")
    mirr = ctx.get("mirr_projeto")
    wacc = ctx.get("wacc")
    pi = ctx.get("pi_projeto")

    cards_r2 = [
        ("IRR", f"{irr:.1f}%" if irr is not None else "N/A",
         irr is not None and irr > ctx["taxa_desc"]),
        ("MIRR", f"{mirr:.1f}%" if mirr is not None else "N/A",
         mirr is not None and mirr > (wacc or 0)),
        ("WACC", f"{wacc:.2f}%" if wacc else "N/A", None),
        ("PI", f"{pi:.2f}x" if pi is not None else "N/A",
         pi is not None and pi > 1),
    ]
    for i, (lbl, val, pos) in enumerate(cards_r2):
        _add_metric_card(slide, x0 + i * (card_w + gap), y2, card_w, card_h, lbl, val, pos)

    # Context line
    ctx_text = (f"{L['slide_sector']}: {ctx['setor']}  |  "
                f"{L['slide_horizon']}: {horizonte} {L['slide_months']}  |  "
                f"CapEx: {fv(ctx['total_capex'], unit)}  |  "
                f"{L['slide_debt']}: {fv(ctx['total_debt'], unit)}  |  "
                f"Equity: {fv(ctx['equity_required'], unit)}")
    _add_text_box(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
                  ctx_text, font_size=10, color=GRAY)

    # Description
    if ctx.get("descricao"):
        _add_text_box(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(0.8),
                      ctx["descricao"], font_size=11, color=BLACK)


def _slide_assumptions(prs, ctx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_assumptions_title"])
    L = ctx["L"]
    fv = ctx["fv"]
    unit = ctx["unit"]

    y = Inches(1.3)
    # Revenue lines
    _add_text_box(slide, Inches(0.6), y, Inches(12), Inches(0.3),
                  L["slide_revenue_lines"], font_size=14, bold=True, color=BLUE_DARK)
    y += Inches(0.35)

    rec_data = [[L["slide_name"], L["slide_annual_value"], L["slide_growth"], L["slide_index"]]]
    for r in ctx["receitas"]:
        rec_data.append([
            r["nome"],
            fv(r["valor_anual"], unit),
            f"{r['crescimento']:.0f}%",
            r["idx"]
        ])
    _add_table(slide, Inches(0.6), y, Inches(8), len(rec_data), 4, rec_data)
    y += Inches(0.35 * len(rec_data) + 0.3)

    # Key parameters
    _add_text_box(slide, Inches(0.6), y, Inches(12), Inches(0.3),
                  L["slide_key_params"], font_size=14, bold=True, color=BLUE_DARK)
    y += Inches(0.35)

    params = [
        [L["slide_param"], L["slide_value"]],
        [L["slide_horizon"], f"{ctx['horizonte']} {L['slide_months']}"],
        [L["slide_disc_rate"], f"{ctx['taxa_desc']:.2f}%"],
        [f"CapEx {L['slide_total']}", fv(ctx["total_capex"], unit)],
        [L["slide_debt"], fv(ctx["total_debt"], unit)],
        ["Equity", fv(ctx["equity_required"], unit)],
        [L["slide_tax_regime"], ctx.get("regime_fiscal", "Lucro Real")],
    ]
    _add_table(slide, Inches(0.6), y, Inches(5), len(params), 2, params)


def _slide_cash_flow_chart(prs, ctx):
    """Cumulative FCF chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_cashflow_title"],
                  ctx["L"]["slide_cashflow_sub"])

    fig = ctx.get("fig_cashflow")
    if fig:
        fig.update_layout(
            width=1100, height=480,
            margin=dict(t=30, b=50, l=60, r=30),
            paper_bgcolor="white", plot_bgcolor="white")
        _add_chart_image(slide, fig, Inches(0.6), Inches(1.3),
                         Inches(12), Inches(5.5))


def _slide_revenue_cost_chart(prs, ctx):
    """Revenue vs Costs chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_rev_cost_title"])

    fig = ctx.get("fig_revcost")
    if fig:
        fig.update_layout(
            width=1100, height=480,
            margin=dict(t=30, b=50, l=60, r=30),
            paper_bgcolor="white", plot_bgcolor="white")
        _add_chart_image(slide, fig, Inches(0.6), Inches(1.3),
                         Inches(12), Inches(5.5))


def _slide_dre(prs, ctx):
    """Income Statement slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_dre_title"],
                  ctx["L"]["slide_dre_sub"].format(unit=ctx["unit"]))

    annual = ctx["annual"]
    anos = list(annual.keys())
    umult = ctx["umult"]

    def mn(v): return f"{v / umult:,.2f}"
    def pn(v): return f"{v:.1f}%"

    headers = [""] + [yr.replace("Ano", "Year") if ctx["lang"] == "EN" else yr for yr in anos]
    L = ctx["L"]
    row_defs = [
        (L["slide_dre_revenue"],     lambda d: mn(d["receita"])),
        (L["slide_dre_cogs"],        lambda d: mn(d["cpv"])),
        (L["slide_dre_gross"],       lambda d: mn(d["lb"])),
        (L["slide_dre_gross_m"],     lambda d: pn(d["mb_pct"])),
        ("EBITDA",                   lambda d: mn(d["ebitda"])),
        (L["slide_dre_ebitda_m"],    lambda d: pn(d["ebitda_pct"])),
        ("EBIT",                     lambda d: mn(d["ebit"])),
        (L["slide_dre_interest"],    lambda d: mn(d["juros"])),
        (L["slide_dre_net"],         lambda d: mn(d["ni"])),
        (L["slide_dre_net_m"],       lambda d: pn(d["ni_pct"])),
    ]
    data = [headers]
    for label, fn in row_defs:
        row = [label] + [fn(annual[yr]) for yr in anos]
        data.append(row)

    _add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5),
               len(data), len(headers), data)


def _slide_debt_schedule(prs, ctx):
    """Debt schedule chart slide."""
    if ctx["total_debt"] <= 0:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_debt_title"])

    fig = ctx.get("fig_debt")
    if fig:
        fig.update_layout(
            width=1100, height=480,
            margin=dict(t=30, b=50, l=60, r=30),
            paper_bgcolor="white", plot_bgcolor="white")
        _add_chart_image(slide, fig, Inches(0.6), Inches(1.3),
                         Inches(12), Inches(5.5))


def _slide_sensitivity(prs, ctx):
    """Sensitivity / tornado chart slide."""
    fig = ctx.get("fig_tornado")
    if not fig:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_sens_title"])

    fig.update_layout(
        width=1100, height=480,
        margin=dict(t=30, b=50, l=60, r=30),
        paper_bgcolor="white", plot_bgcolor="white")
    _add_chart_image(slide, fig, Inches(0.6), Inches(1.3),
                     Inches(12), Inches(5.5))


def _slide_appendix(prs, ctx):
    """Monthly table (first 24 months max)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_header(slide, ctx["L"]["slide_appendix_title"],
                  ctx["L"]["slide_appendix_sub"])

    df_op = ctx["df_op"]
    umult = ctx["umult"]
    max_months = min(24, len(df_op))
    df_show = df_op.head(max_months)

    headers = [ctx["L"]["slide_month"], ctx["L"]["slide_dre_revenue"],
               "CPV", "EBIT", "FCF",
               ctx["L"]["slide_accumulated"]]
    data = [headers]
    for _, row in df_show.iterrows():
        data.append([
            str(int(row["Mes"])),
            f"{row['Receita']/umult:,.1f}",
            f"{row['CPV']/umult:,.1f}",
            f"{row['EBIT']/umult:,.1f}",
            f"{row['FCF']/umult:,.1f}",
            f"{row['Acumulado']/umult:,.1f}",
        ])

    if len(data) > 1:
        _add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5),
                   len(data), len(headers), data)


# ─── Main generator ──────────────────────────────────────────────────────────

def gerar_deck(ctx):
    """Generate a PPTX deck and return bytes.

    ctx: dict with all model data and translation keys.
    Required keys: nome_proj, setor, responsavel, data_inicio, descricao,
        receitas, horizonte, taxa_desc, total_capex, total_debt, equity_required,
        npv, payback, roi, mg_med, irr_projeto, mirr_projeto, wacc, pi_projeto,
        annual, df_op, unit, umult, fv, lang, L,
        fig_cashflow, fig_revcost, fig_debt (optional Plotly figs),
        fig_tornado (optional).
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, ctx)
    _slide_executive_summary(prs, ctx)
    _slide_assumptions(prs, ctx)
    _slide_cash_flow_chart(prs, ctx)
    _slide_revenue_cost_chart(prs, ctx)
    _slide_dre(prs, ctx)
    _slide_debt_schedule(prs, ctx)
    _slide_sensitivity(prs, ctx)
    _slide_appendix(prs, ctx)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─── Translation keys for slides ────────────────────────────────────────────

SLIDE_TRANSLATIONS = {
    "PT": {
        "slide_cover_subtitle": "Analise de Viabilidade Financeira",
        "slide_sector": "Setor", "slide_author": "Autor", "slide_date": "Data",
        "slide_exec_title": "Sumario Executivo",
        "slide_exec_sub": "Visao geral do projeto e metricas de viabilidade",
        "slide_viable": "PROJETO VIAVEL", "slide_caution": "ATENCAO — VIABILIDADE PARCIAL",
        "slide_not_viable": "PROJETO INVIAVEL",
        "slide_roi": "ROI", "slide_npv": "VPL", "slide_payback": "Payback",
        "slide_margin": "Margem Bruta", "slide_months": "meses",
        "slide_horizon": "Horizonte", "slide_debt": "Divida", "slide_total": "total",
        "slide_disc_rate": "Taxa de desconto", "slide_tax_regime": "Regime tributario",
        "slide_assumptions_title": "Premissas do Modelo",
        "slide_revenue_lines": "Linhas de Receita",
        "slide_name": "Nome", "slide_annual_value": "Valor anual",
        "slide_growth": "Crescimento", "slide_index": "Indexador",
        "slide_key_params": "Parametros-chave",
        "slide_param": "Parametro", "slide_value": "Valor",
        "slide_cashflow_title": "Fluxo de Caixa Acumulado",
        "slide_cashflow_sub": "Projeto (unlevered) vs. Equity (levered)",
        "slide_rev_cost_title": "Receita vs. Custos",
        "slide_dre_title": "Demonstracao de Resultado (DRE)",
        "slide_dre_sub": "Valores em {unit}",
        "slide_dre_revenue": "Receita Bruta", "slide_dre_cogs": "(–) CPV",
        "slide_dre_gross": "Lucro Bruto", "slide_dre_gross_m": "Margem Bruta (%)",
        "slide_dre_ebitda_m": "Margem EBITDA (%)",
        "slide_dre_interest": "(–) Desp. Financeiras",
        "slide_dre_net": "Lucro Liquido", "slide_dre_net_m": "Margem Liquida (%)",
        "slide_debt_title": "Estrutura de Divida — Schedule Consolidado",
        "slide_sens_title": "Analise de Sensibilidade — Tornado",
        "slide_appendix_title": "Apendice — Tabela Mensal",
        "slide_appendix_sub": "Primeiros 24 meses (valores em unidade do modelo)",
        "slide_month": "Mes", "slide_accumulated": "Acumulado",
        "slide_download": "Baixar deck (.pptx)",
        "slide_generating": "Gerando deck...",
        "slide_ready": "Deck gerado com sucesso!",
        "slide_filename": "business_case_{nome}.pptx",
    },
    "EN": {
        "slide_cover_subtitle": "Financial Viability Analysis",
        "slide_sector": "Sector", "slide_author": "Author", "slide_date": "Date",
        "slide_exec_title": "Executive Summary",
        "slide_exec_sub": "Project overview and viability metrics",
        "slide_viable": "VIABLE PROJECT", "slide_caution": "CAUTION — PARTIAL VIABILITY",
        "slide_not_viable": "NOT VIABLE PROJECT",
        "slide_roi": "ROI", "slide_npv": "NPV", "slide_payback": "Payback",
        "slide_margin": "Gross Margin", "slide_months": "months",
        "slide_horizon": "Horizon", "slide_debt": "Debt", "slide_total": "total",
        "slide_disc_rate": "Discount rate", "slide_tax_regime": "Tax regime",
        "slide_assumptions_title": "Model Assumptions",
        "slide_revenue_lines": "Revenue Lines",
        "slide_name": "Name", "slide_annual_value": "Annual value",
        "slide_growth": "Growth", "slide_index": "Index",
        "slide_key_params": "Key Parameters",
        "slide_param": "Parameter", "slide_value": "Value",
        "slide_cashflow_title": "Cumulative Cash Flow",
        "slide_cashflow_sub": "Project (unlevered) vs. Equity (levered)",
        "slide_rev_cost_title": "Revenue vs. Costs",
        "slide_dre_title": "Income Statement",
        "slide_dre_sub": "Values in {unit}",
        "slide_dre_revenue": "Gross Revenue", "slide_dre_cogs": "(–) COGS",
        "slide_dre_gross": "Gross Profit", "slide_dre_gross_m": "Gross Margin (%)",
        "slide_dre_ebitda_m": "EBITDA Margin (%)",
        "slide_dre_interest": "(–) Interest Expense",
        "slide_dre_net": "Net Income", "slide_dre_net_m": "Net Margin (%)",
        "slide_debt_title": "Debt Structure — Consolidated Schedule",
        "slide_sens_title": "Sensitivity Analysis — Tornado",
        "slide_appendix_title": "Appendix — Monthly Table",
        "slide_appendix_sub": "First 24 months (values in model unit)",
        "slide_month": "Month", "slide_accumulated": "Accumulated",
        "slide_download": "Download deck (.pptx)",
        "slide_generating": "Generating deck...",
        "slide_ready": "Deck generated successfully!",
        "slide_filename": "business_case_{nome}.pptx",
    },
}
